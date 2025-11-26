# create_presentation.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import imageio
from PIL import Image, ImageDraw, ImageFont
import os

# --- Professional Theme Colors ---
THEME_COLOR_BLUE = RGBColor(0, 82, 165)
THEME_COLOR_GRAY = RGBColor(128, 128, 128)
COLOR_CORTEX_ZONE = RGBColor(220, 230, 242) # Light Blue
COLOR_TARGET_ZONE = RGBColor(226, 240, 217) # Light Green
COLOR_SHARED_ZONE = RGBColor(242, 242, 242) # Light Gray
COLOR_INTERACTION = RGBColor(255, 242, 204) # Light Yellow
COLOR_ARROW = RGBColor(191, 144, 0) # Gold-ish
COLOR_HIGHLIGHT = RGBColor(255, 25, 25) # Red



# --- Presentation Content ---
# Each dictionary in this list represents a slide.
slides_content = [
    {
        "title": "Cortex Performance Engine",
        "subtitle": "An Intelligent, Agentic Resilience & Performance Platform on AWS",
        "notes": "The title slide sets the stage for the entire presentation. It should be clean, professional, and clearly state the project's name and its high-level purpose."
    },
    {
        "title": "Executive Summary: The Opportunity",
        "points": [
            {"text": "The Problem:", "bold": True, "level": 0, "end_para": " Our current performance testing is slow, manual, and costly, creating a bottleneck for innovation and leaving us vulnerable to production failures."},
            {"text": "The Solution:", "bold": True, "level": 0, "end_para": " The Cortex Performance Engine is a serverless, AI-driven platform that automates the entire resilience engineering lifecycle, from test creation to root cause analysis."},
            {"text": "The Ask:", "bold": True, "level": 0, "end_para": " We seek approval to deploy the prototype and integrate it with our primary e-commerce application to validate its business impact."},
            {"text": "The ROI:", "bold": True, "level": 0, "end_para": " We project a 90%+ reduction in testing infrastructure costs and a 95%+ reduction in manual effort, enabling faster delivery and more resilient applications."}
        ],
        "notes": "This slide is crucial for stakeholders. It immediately answers 'What is this about and why should I care?'"
    },
    {
        "title": "Agenda",
        "points": [
            "1. The Strategic Objective: Why We Need This",
            "2. The Problem with the Status Quo",
            "3. Our Solution: The Cortex Performance Engine",
            "4. Architecture Deep Dive",
            "5. Key Advantages & Business Impact",
            "6. Live Prototype Demonstration",
            "7. Financials: Cost & ROI Analysis",
            "8. Roadmap & Next Steps"
        ]
    },
    {
        "title": "The Strategic Objective: Speed with Stability",
        "points": [
            "Our core business objective is to accelerate feature delivery while ensuring rock-solid application stability and performance.",
            "To achieve this, we must move from a reactive, manual testing model to a proactive, automated resilience validation culture.",
            {"text": "The Cortex Performance Engine is the key enabler of this strategic shift.", "level": 1, "italic": True}
        ]
    },
    {
        "title": "The Problem: The High Cost of Manual Testing",
        "subtitle": "Our current approach is a major bottleneck, characterized by high costs, slow feedback, and significant risk.",
        "table": {
            "headers": ["Challenge", "Description", "Impact"],
            "rows": [
                ["High Manual Effort", "Engineers spend weeks writing brittle test scripts.", "Slows down release cycles; High personnel cost."],
                ["High Infrastructure Cost", "Requires dedicated, always-on servers for test controllers.", "Wasteful spending; Significant maintenance overhead."],
                ["Disconnected Data", "Test results, server metrics, and logs are in different systems.", "Difficult to correlate cause and effect; Slow root cause analysis."],
                ["Production Gaps", "Tests are based on assumptions, not real user behavior.", "False sense of security; Production incidents still occur."]
            ]
        }
    },
    {
        "title": "The Solution: Cortex Performance Engine",
        "subtitle": "A serverless, AI-driven platform that automates the entire resilience engineering lifecycle.",
        "points": [
            {"text": "Chatbot-Driven:", "bold": True, "level": 0, "end_para": " Users define and launch complex tests using simple, natural language."},
            {"text": "AI-Powered:", "bold": True, "level": 0, "end_para": " Generative AI (Amazon Bedrock) analyzes logs, generates test scripts, and writes expert-level performance reports automatically."},
            {"text": "Fully Automated:", "bold": True, "level": 0, "end_para": " The platform is orchestrated by AWS Step Functions, requiring zero manual intervention from start to finish."},
            {"text": "Cost-Optimized:", "bold": True, "level": 0, "end_para": " Built on a serverless foundation (Lambda, Fargate Spot) to minimize costs, projecting a 90%+ reduction in test infrastructure spend."}
        ],
        "notes": "This slide clearly articulates the four key pillars of the solution."
    },
    {
        "title": "Architecture 1: The E-Commerce Web Application",
        "subtitle": "A standard, serverless 3-Tier architecture that serves as the target for our tests.",
        "is_diagram": True,
        "gif_filename": "ecommerce_architecture.gif",
        "diagram_elements": [
            # --- Groups ---
            {"id": "group1", "text": "Presentation Tier", "x": 0.5, "y": 2, "w": 15, "h": 2.5, "is_group": True, "font_size": 14, "color": COLOR_TARGET_ZONE},
            {"id": "group2", "text": "Application Tier", "x": 0.5, "y": 4.75, "w": 15, "h": 2, "is_group": True, "font_size": 14, "color": COLOR_TARGET_ZONE},
            {"id": "group3", "text": "Data Tier", "x": 0.5, "y": 7, "w": 15, "h": 1.5, "is_group": True, "font_size": 14, "color": COLOR_TARGET_ZONE},
            # Components
            {"id": "user", "text": "End User\n(Browser)", "x": 1, "y": 3, "w": 2, "h": 1, "color": COLOR_INTERACTION},
            {"id": "cdn", "text": "Amazon CloudFront\n(CDN)", "x": 4, "y": 3, "w": 3, "h": 1},
            {"id": "s3", "text": "S3 Bucket\n(React Static Files)", "x": 8, "y": 3, "w": 3, "h": 1},
            {"id": "apigw", "text": "Amazon API Gateway\n(REST API)", "x": 4, "y": 5.25, "w": 3, "h": 1},
            {"id": "lambda", "text": "AWS Lambda\n(FastAPI Backend Logic)", "x": 8, "y": 5.25, "w": 3, "h": 1},
            {"id": "db", "text": "Amazon DynamoDB\n(Database)", "x": 8, "y": 7.25, "w": 3, "h": 1},
            # Arrows
            {"id": "arrow1", "text": "HTTPS Request", "x": 2.5, "y": 2.5, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow2", "text": "Serves Frontend", "x": 6.5, "y": 2.5, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow3", "text": "API Calls (e.g., /api/products)", "x": 5, "y": 4.25, "w": 3, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow4", "text": "Invokes", "x": 6.5, "y": 6.4, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow5", "text": "Reads/Writes Data", "x": 9, "y": 6.75, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
        ],
        "animation_flow": ["user", "cdn", "s3", "apigw", "lambda", "db"],
        "notes": "This is our target application. It's a modern, serverless web app. Understanding its components is key to understanding how we test it."
    },
    {
        "title": "Architecture 2: The Cortex Performance Engine",
        "subtitle": "An event-driven, microservice-based platform for intelligent testing.",
        "is_diagram": True,
        "gif_filename": "cortex_engine_architecture.gif",
        "diagram_elements": [
            # --- Groups ---
            {"id": "group1", "text": "Interaction Layer", "x": 0.5, "y": 2, "w": 15, "h": 1.75, "is_group": True, "font_size": 14, "color": COLOR_CORTEX_ZONE},
            {"id": "group2", "text": "Orchestration Layer", "x": 0.5, "y": 4, "w": 15, "h": 1.75, "is_group": True, "font_size": 14, "color": COLOR_CORTEX_ZONE},
            {"id": "group3", "text": "Agentic Core (Execution Layer)", "x": 0.5, "y": 6, "w": 15, "h": 2.5, "is_group": True, "font_size": 14, "color": COLOR_CORTEX_ZONE},
            # Components
            {"id": "user", "text": "User", "x": 1, "y": 2.5, "w": 2, "h": 1, "color": COLOR_INTERACTION},
            {"id": "lex", "text": "AWS Lex\n(Chatbot)", "x": 4, "y": 2.5, "w": 3, "h": 1},
            {"id": "trigger", "text": "Trigger Lambda", "x": 8, "y": 2.5, "w": 3, "h": 1},
            {"id": "sfn", "text": "AWS Step Functions\n(State Machine)", "x": 4, "y": 4.5, "w": 8, "h": 1},
            {"id": "agent1", "text": "Log Analyzer\n(Lambda)", "x": 1, "y": 6.75, "w": 2.5, "h": 1},
            {"id": "agent2", "text": "Script Generator\n(Bedrock)", "x": 4, "y": 6.75, "w": 2.5, "h": 1},
            {"id": "agent3", "text": "Test Executor\n(Fargate Spot)", "x": 7, "y": 6.75, "w": 2.5, "h": 1},
            {"id": "agent4", "text": "Report Synthesizer\n(Bedrock)", "x": 10, "y": 6.75, "w": 2.5, "h": 1},
            {"id": "agent5", "text": "n8n Webhook\n(Notifications)", "x": 13, "y": 6.75, "w": 2.5, "h": 1},
            # Arrows
            {"id": "arrow1", "text": "Natural Language", "x": 2.5, "y": 2, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow2", "text": "Parses Intent", "x": 6.5, "y": 2, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow3", "text": "Starts Workflow", "x": 7.5, "y": 3.6, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow4", "text": "Orchestrates Agents in Sequence", "x": 7.5, "y": 5.6, "w": 3, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
        ],
        "animation_flow": ["user", "lex", "trigger", "sfn", "agent1", "agent2", "agent3", "agent4", "agent5"],
        "notes": "This is our solution. It's an event-driven system where each agent is a specialized microservice. Step Functions acts as the central brain, ensuring tasks run in the correct order."
    },
    {
        "title": "Architecture 3: AWS Services Rationale",
        "subtitle": "Choosing the right tool for the job to maximize efficiency and minimize cost.",
        "table": {
            "headers": ["Category", "Service Chosen", "Why We Chose It (The Rationale)"],
            "rows": [
                ["User Interface", "AWS Lex", "Provides a natural language interface, making the powerful backend accessible to everyone without a complex UI."],
                ["Orchestration", "AWS Step Functions", "Gives us visual workflows, built-in error handling, and state management. It's the perfect serverless 'brain' for our agentic system."],
                ["AI / Machine Learning", "Amazon Bedrock", "Provides secure, serverless access to powerful foundation models for script generation and report analysis, without any ML overhead."],
                ["Serverless Compute", "AWS Lambda", "Ideal for short-lived, event-driven tasks like our agents. We only pay per millisecond, ensuring maximum cost-efficiency."],
                ["Container Compute", "AWS Fargate Spot", "The best of both worlds: serverless container orchestration with massive (up to 90%) cost savings. Perfect for the JMeter load generators."]
            ]
        },
        "notes": "Our architectural choices are deliberate, focusing on serverless-first principles to reduce operational overhead and cost at every layer."
    },
    {
        "title": "Consolidated Architecture: The Full Ecosystem",
        "subtitle": "How the Cortex Engine applies intelligent testing to the E-Commerce Application.",
        "is_diagram": True,
        "gif_filename": "consolidated_architecture.gif",
        "diagram_elements": [
            # --- Groups ---
            {"id": "group1", "text": "Cortex Performance Engine", "x": 0.5, "y": 1.5, "w": 7.5, "h": 6.5, "is_group": True, "font_size": 14, "color": COLOR_CORTEX_ZONE},
            {"id": "group2", "text": "Target E-Commerce Application", "x": 8.25, "y": 1.5, "w": 7.25, "h": 6.5, "is_group": True, "font_size": 14, "color": COLOR_TARGET_ZONE},
            {"id": "group3", "text": "Shared Services & Data Flow", "x": 0.5, "y": 8.25, "w": 15, "h": 0.5, "is_group": True, "font_size": 14, "color": COLOR_SHARED_ZONE},
            # --- Cortex Zone ---
            {"id": "user", "text": "User", "x": 1, "y": 2, "w": 1.5, "h": 0.75, "color": COLOR_INTERACTION},
            {"id": "lex", "text": "AWS Lex", "x": 4, "y": 2, "w": 1.5, "h": 0.75, "color": COLOR_INTERACTION},
            {"id": "sfn", "text": "Step Functions", "x": 1, "y": 4.5, "w": 3, "h": 1.5},
            {"id": "agents", "text": "Agent Core", "x": 4.5, "y": 4, "w": 3, "h": 3.5},
            # --- Target Zone ---
            {"id": "apigw", "text": "API Gateway", "x": 8.75, "y": 3.5, "w": 2.5, "h": 1},
            {"id": "lambda", "text": "Lambda Backend", "x": 12.25, "y": 3.5, "w": 2.75, "h": 1},
            # --- Shared Zone ---
            {"id": "s3", "text": "S3 Artifacts Bucket", "x": 1, "y": 8.25, "w": 6.5, "h": 0.5},
            {"id": "cw", "text": "CloudWatch Logs", "x": 8.5, "y": 8.25, "w": 6.5, "h": 0.5},
            # --- Arrows ---
            {"id": "arrow1", "text": "1. 'Run a test...'", "x": 2.5, "y": 2.15, "w": 1.5, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow2", "text": "2. Starts Workflow", "x": 2, "y": 3.5, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow3", "text": "3. Reads Logs", "x": 7.5, "y": 7.25, "w": 2, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow4", "text": "4. Writes/Reads Scripts & Reports", "x": 1, "y": 7.25, "w": 3.5, "h": 0.5, "font_size": 11, "no_box": True, "color": COLOR_ARROW},
            {"id": "arrow5", "text": "5. APPLIES LOAD", "x": 7.5, "y": 2.5, "w": 2, "h": 0.5, "font_size": 12, "no_box": True, "color": RGBColor(255, 0, 0), "bold": True},
        ],
        "animation_flow": ["user", "lex", "sfn", "agents", "cw", "s3", "apigw", "lambda"],
        "notes": "This is the complete picture. The Cortex Engine (left) is a fully independent system that interacts with the Target Application (right) only through public APIs and by observing logs, just like a real user and an SRE would."
    },
    {
        "title": "Key Advantages: The Business Impact",
        "points": [
            {"text": "Accelerated Delivery:", "bold": True, "level": 0, "end_para": " Reduce test creation time from weeks to minutes, removing bottlenecks and enabling faster time-to-market."},
            {"text": "Drastic Cost Reduction:", "bold": True, "level": 0, "end_para": " Save over 90% on infrastructure costs by leveraging a serverless, pay-per-use model instead of always-on servers."},
            {"text": "Increased Resilience:", "bold": True, "level": 0, "end_para": " Proactively find and fix performance and chaos issues before they impact customers, reducing the risk of costly outages."},
            {"text": "Improved Developer Productivity:", "bold": True, "level": 0, "end_para": " Automate tedious tasks, freeing up expert engineers to focus on innovation and high-value work."},
            {"text": "Democratized Testing:", "bold": True, "level": 0, "end_para": " Empower more team members to run sophisticated tests via a simple chatbot, fostering a culture of quality."}
        ]
    },
    {
        "title": "Live Prototype Demonstration",
        "points": [
            {"text": "1. Triggering the Test:", "bold": True, "level": 0, "end_para": " We will interact with the chatbot in the AWS Lex console to launch a load test."},
            {"text": "2. Monitoring the Pipeline:", "bold": True, "level": 0, "end_para": " We will observe the AWS Step Functions graph to see the automated workflow in real-time."},
            {"text": "3. Viewing the AI-Generated Report:", "bold": True, "level": 0, "end_para": " We will review the final PDF report in S3, complete with metrics, analysis, and AI-generated recommendations."}
        ],
        "notes": "This slide sets the agenda for the live demo portion of the presentation."
    },
    {
        "title": "Financials: Cost & ROI Analysis",
        "subtitle": "A compelling business case built on dramatic cost reduction and risk mitigation.",
        "table": {
            "headers": ["Category", "Traditional Manual Approach", "Agentic Platform", "Savings"],
            "rows": [
                ["Personnel Effort (per test)", "40-80 hours (~$4k-$8k)", "< 1 hour (~$100)", ">98% Reduction"],
                ["Test Infrastructure (monthly)", "~$500+ (always-on servers)", "<$30 (pay-per-use)", ">94% Reduction"],
                ["Reporting & Analysis (per test)", "4-8 hours (~$400-$800)", "0 hours (fully automated)", "100% Reduction"],
            ]
        },
        "notes": "The ROI is clear. The platform pays for itself by avoiding the cost of a single engineer's time for one week, while continuously reducing infrastructure spend and mitigating outage risks."
    },
    {
        "title": "Roadmap & Next Steps",
        "points": [
            {"text": "Phase 1: Deploy & Validate (This Quarter)", "bold": True, "level": 0},
            {"text": "Deploy the Cortex Engine prototype to AWS.", "level": 1},
            {"text": "Integrate with the e-commerce application.", "level": 1},
            {"text": "Execute baseline tests and validate cost savings and performance metrics.", "level": 1},
            {"text": "Phase 2: Enhance & Expand (Next Quarter)", "bold": True, "level": 0},
            {"text": "Onboard additional development teams.", "level": 1},
            {"text": "Enhance the AI reporting agent with more advanced diagnostics.", "level": 1},
            {"text": "Phase 3: Autonomous Operation (Future)", "bold": True, "level": 0},
            {"text": "Implement proactive, scheduled resilience checks.", "level": 1},
            {"text": "Explore self-healing capabilities and automated remediation suggestions.", "level": 1}
        ]
    },
    {
        "title": "Thank You & Q&A",
        "subtitle": "Opening the floor for questions."
    }
]

def add_footer(slide, slide_num, total_slides):
    """Adds a footer with slide number and project name to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(15), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Cortex Performance Engine | {slide_num} of {total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = THEME_COLOR_GRAY
    p.alignment = PP_ALIGN.RIGHT

def create_presentation(filename="Cortex_Performance_Engine_Stakeholder_Presentation.pptx"):
    """
    Generates a PowerPoint presentation from the slides_content list.
    """
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    total_slides = len(slides_content)

    for i, slide_data in enumerate(slides_content):
        is_title_slide = (i == 0)
        
        if is_title_slide:
            slide_layout = prs.slide_layouts[0] # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide_layout = prs.slide_layouts[5] # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data.get("title", "")
        title.text_frame.paragraphs[0].font.color.rgb = THEME_COLOR_BLUE
        title.text_frame.paragraphs[0].font.bold = True

        body_shape = None
        # Find the main content placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break

        if "subtitle" in slide_data and body_shape:
            body_shape.text = slide_data.get("subtitle", "")
            body_shape.text_frame.paragraphs[0].font.size = Pt(24)
            body_shape.text_frame.paragraphs[0].font.italic = True
            body_shape.text_frame.paragraphs[0].font.color.rgb = THEME_COLOR_GRAY

        if "points" in slide_data:
            tf = body_shape.text_frame if body_shape else slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(14), Inches(6.5)).text_frame
            tf.clear()
            tf.word_wrap = True
            
            for point in slide_data["points"]:
                if isinstance(point, dict):
                    p = tf.add_paragraph()
                    p.text = point.get("text", "")
                    p.level = point.get("level", 0)
                    p.font.bold = point.get("bold", False)
                    p.font.italic = point.get("italic", False)
                    p.font.size = Pt(24)
                    if "end_para" in point:
                        run = p.add_run()
                        run.text = point["end_para"]
                        run.font.bold = False
                        run.font.italic = False
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(24)

        if "table" in slide_data:
            table_data = slide_data["table"]
            rows, cols = len(table_data["rows"]) + 1, len(table_data["headers"])
            shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.5), Inches(14), Inches(3))
            table = shape.table

            for c, header in enumerate(table_data["headers"]):
                table.cell(0, c).text = header
                table.cell(0, c).text_frame.paragraphs[0].font.bold = True
                table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(16)
                table.cell(0, c).text_frame.paragraphs[0].font.color.rgb = THEME_COLOR_BLUE
                table.cell(0, c).fill.solid()
                table.cell(0, c).fill.fore_color.rgb = COLOR_SHARED_ZONE

            for r, row_data in enumerate(table_data["rows"]):
                for c, cell_text in enumerate(row_data):
                    table.cell(r + 1, c).text = cell_text
                    table.cell(r + 1, c).text_frame.paragraphs[0].font.size = Pt(14)

        if slide_data.get("is_diagram"):
            gif_filename = slide_data.get("gif_filename")
            if gif_filename:
                # Generate the GIF
                generate_animated_diagram(
                    gif_filename,
                    slide_data["diagram_elements"],
                    slide_data.get("animation_flow", [])
                )
                # Add the GIF to the slide
                slide.shapes.add_picture(gif_filename, Inches(0.5), Inches(1.8), width=Inches(15))
        
        if "notes" in slide_data:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["notes"]

        if not is_title_slide:
            add_footer(slide, i + 1, total_slides)

    print(f"Presentation '{filename}' created successfully.")
    prs.save(filename)

if __name__ == '__main__':

    create_presentation()
